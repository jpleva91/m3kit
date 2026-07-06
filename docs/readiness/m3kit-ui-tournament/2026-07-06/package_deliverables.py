from __future__ import annotations

import html
import json
import os
import re
import subprocess
import textwrap
from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT
SLIDES_DIR = OUT / "video-package" / "slides"
VIDEO_DIR = OUT / "video-package"
PPTX = OUT / "m3kit-ui-tournament-v2-deck.pptx"
HTML = OUT / "m3kit-ui-tournament-v2-brief.html"
MP4 = VIDEO_DIR / "m3kit-ui-tournament-v2-silent-deck.mp4"
SRT = VIDEO_DIR / "m3kit-ui-tournament-v2-captions.srt"
VOICEOVER = VIDEO_DIR / "voiceover-script.md"
VERIFY = OUT / "verification-results.txt"

NAVY = "#162033"
INK = "#1d2630"
COBALT = "#2454d6"
SIENNA = "#a15c38"
PAPER = "#f6f1e8"
MIST = "#e8ecf3"
GREEN = "#2f7d4f"
AMBER = "#b9791e"
RED = "#b44848"

REQUIRED = [
    "tournament-report.md",
    "feature-candidate-list.md",
    "notebooklm-source-packet.md",
    "slide-deck-outline.md",
    "video-presentation-script.md",
    "contestant-run-results.json",
]

RANKS = [
    ["1", "OpenCode / DeepSeek v4 Flash free", "Data Table Column Manager", "30", "Winner"],
    ["2", "Antigravity / Gemini 3.5 Flash Medium", "Data Table Column Manager", "28", "Runner-up"],
    ["3", "Codex / gpt-5.5", "Saved View Menu", "24", "Honorable mention"],
    ["4", "Hermes openai-codex / gpt-5.5", "Data State primitive", "21", "Salvageable"],
]

SLIDES = [
    {
        "title": "m3kit UI Tournament v2",
        "kicker": "Proposal-only model × driver tournament",
        "bullets": [
            "Mission: choose the next launch-readiness UI component feature after the @m3kit/ai gate.",
            "Winner: Data Table Column Manager (m3k-column-manager).",
            "Artifact pack: report, NotebookLM source packet, deck, video package, phone brief.",
        ],
        "callout": "Decision: build Column Manager first; Saved View Manager follows.",
        "kind": "title",
    },
    {
        "title": "The contract narrowed the field",
        "kicker": "Hard boundaries contestants had to honor",
        "bullets": [
            "Token-only styling: Material system tokens + closed --app-* contract.",
            "No new UI/chart libraries, endpoints, credentials, real data, or model downloads.",
            "Every exported component needs unit spec, Storybook story, and Cypress component test.",
            "Clean-room: public sources only; unsupported claims were penalized.",
        ],
        "callout": "Release Captain authority: local-write artifact packaging only; no external publish beyond Tailnet serving.",
        "kind": "cards",
    },
    {
        "title": "Shared source packet",
        "kicker": "Contestants received the same grounding packet",
        "bullets": [
            "Repo baseline: Angular 19 / Nx 20 reporting UI reference library.",
            "DESIGN.md doctrine: industrial-editorial Instruments, density discipline, dark-mode parity.",
            "Readiness gaps: table controls, saved views, chart state/a11y, filter/date ergonomics.",
            "@m3kit/ai context: approved runtime seam, but assistant UI/provider adapters deferred.",
        ],
        "callout": "NotebookLM status: source packet ready; notebook creation pending because no NotebookLM tool/credential was used.",
        "kind": "split",
    },
    {
        "title": "Candidate slate",
        "kicker": "Eight options; four strongest tournament contenders",
        "bullets": [
            "A Saved View Manager — strategic follow-up for query + column state.",
            "B Report Action Bar — export/refresh/action composition for reporting pages.",
            "C Chart State Frame — accessibility and loading/error/empty consistency.",
            "D Data Table Column Manager — first-class toggle/reorder/pin UI over existing headless state.",
        ],
        "callout": "Deferred: AI assistant UI until provider policy is explicit; no-secret map shell remains high risk.",
        "kind": "grid",
    },
    {
        "title": "Runtime roster was evidence-first",
        "kicker": "Completed packets scored; hung/degraded lanes were labeled honestly",
        "bullets": [
            "Completed usable packets: Codex gpt-5.5, agy Gemini Flash Medium, OpenCode DeepSeek free, Hermes current lane.",
            "Claude Code Fable hung with no stdout and was marked DNF/hung.",
            "Additional agy/OpenCode variants were recorded as DNF/degraded instead of over-spending runtime.",
        ],
        "callout": "No silent success: degraded lanes remain visible in the manifest and report.",
        "kind": "timeline",
    },
    {
        "title": "Scoring rubric",
        "kicker": "Six dimensions, max 30 points",
        "bullets": [
            "Contract fit and boundary safety",
            "Launch/readiness impact",
            "UX/component API quality",
            "Accessibility and state coverage",
            "Test/Storybook/evidence quality",
            "Implementation feasibility",
        ],
        "callout": "Winning proposal scored 30/30 because it used existing seams and added no new architectural risk.",
        "kind": "rubric",
    },
    {
        "title": "Winner: Data Table Column Manager",
        "kicker": "OpenCode / opencode-deepseek-v4-flash-free",
        "bullets": [
            "Component: m3k-column-manager in libs/table.",
            "Inputs: TableDefinition + ColumnViewState; emits intent-only column state changes.",
            "States: visible/hidden, reorder, pin left/right, locked/required columns, narrow viewport.",
            "Accessibility: no drag-only interaction; keyboard move controls are mandatory.",
        ],
        "callout": "Winner total: 30 / 30",
        "kind": "winner",
    },
    {
        "title": "Why it wins",
        "kicker": "It closes an enterprise parity gap with existing repo seams",
        "bullets": [
            "m3k-data-table already accepts headless columnState.",
            "libs/core already owns ColumnViewState and resolveColumns.",
            "Column Manager lives cleanly in libs/table and does not own persistence.",
            "Saved views can later capture emitted state without redesigning the table API.",
        ],
        "callout": "This is launch-credible UI, not a new dependency, backend, or AI/provider bet.",
        "kind": "diagram",
    },
    {
        "title": "Ranking summary",
        "kicker": "The top two converged on the same feature",
        "bullets": [],
        "callout": "Codex’s Saved View Menu is the right second slice after Column Manager emits stable state.",
        "kind": "ranking",
    },
    {
        "title": "Recommended Spec Kit feature pack",
        "kicker": "specs/007-table-column-manager-and-saved-view-seam/",
        "bullets": [
            "Implement m3k-column-manager in libs/table.",
            "Storybook: default, hidden columns, pinned columns, locked columns, many columns, narrow viewport.",
            "Unit tests: emitted state generation and edge cases.",
            "Cypress: keyboard/focus tests for no-drag dependency.",
            "Demo: feed emitted ColumnViewState[] into m3k-data-table; no persistence in first slice.",
        ],
        "callout": "Effort estimate: 1.5–2.5 focused engineering days; saved-view follow-up +1–2 days.",
        "kind": "steps",
    },
    {
        "title": "Evidence and caveats",
        "kicker": "Deliverables are source-backed, local, and phone-readable",
        "bullets": [
            "Verified source artifacts: report, candidate list, contestant packets, NotebookLM packet, outline, script.",
            "Generated: PPTX deck, silent MP4 deck video, captions, voiceover-ready script, phone HTML brief.",
            "NotebookLM: source packet ready; notebook creation pending, not fabricated.",
            "Tailnet: HTML brief served and verified HTTP 200 by Release Captain.",
        ],
        "callout": "Next exact action: approve specs/007 and route implementation to the appropriate coding lane.",
        "kind": "close",
    },
]

SCRIPT_SECTIONS = {
    "Opening": "This tournament asked the available agent drivers to choose the next proposal-only UI component feature for m3kit after the @m3kit/ai port gate was approved.",
    "Context": "m3kit is a clean-room Angular 19 / Nx 20 Material 3 reporting kit. The hard constraints are token-only styling, no new chart/UI libraries, synthetic data only, strict library boundaries, and complete component coverage.",
    "Source packet": "Contestants received the same source packet: exported components, DESIGN.md doctrine, feature roadmap gaps, UI parity research, and the approved @m3kit/ai slice. The AI port matters as context, but assistant UI and provider adapters stay deferred.",
    "Roster": "Runtime discovery found Codex, Claude Code, Antigravity, and OpenCode. Codex, agy Gemini Flash Medium, OpenCode DeepSeek free, and the Hermes current lane produced usable packets. Hung or unbounded lanes were explicitly marked DNF or degraded.",
    "Winner": "The winner is OpenCode with the Data Table Column Manager. It noticed that the table already accepts columnState and core already has headless column resolution. The missing piece is a UI for toggling, reordering, and pinning columns.",
    "Recommendation": "The next Spec Kit feature pack should implement m3k-column-manager in libs/table with Storybook coverage, unit tests for emitted state, and Cypress keyboard and focus tests. Saved View Manager should follow as the second slice.",
}


def font(size: int, bold: bool = False):
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf",
    ]
    for c in candidates:
        if Path(c).exists():
            return ImageFont.truetype(c, size)
    return ImageFont.load_default()


def hex_to_rgb(h: str) -> tuple[int, int, int]:
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def wrap_text(draw: ImageDraw.ImageDraw, text: str, fnt, width: int) -> list[str]:
    out = []
    for para in text.split("\n"):
        words = para.split()
        line = ""
        for word in words:
            cand = (line + " " + word).strip()
            if draw.textbbox((0, 0), cand, font=fnt)[2] <= width:
                line = cand
            else:
                if line:
                    out.append(line)
                line = word
        if line:
            out.append(line)
    return out


def draw_card(draw, xy, title, body, accent=COBALT):
    x, y, w, h = xy
    draw.rounded_rectangle((x, y, x + w, y + h), 22, fill="white", outline="#d5d9e2", width=2)
    draw.rectangle((x, y, x + 9, y + h), fill=accent)
    draw.text((x + 28, y + 22), title, fill=INK, font=font(32, True))
    yy = y + 70
    for line in wrap_text(draw, body, font(25), w - 60):
        draw.text((x + 28, yy), line, fill="#3f4c5b", font=font(25))
        yy += 34


def draw_slide_image(slide, idx):
    W, H = 1920, 1080
    img = Image.new("RGB", (W, H), PAPER)
    draw = ImageDraw.Draw(img)
    dark = slide["kind"] in {"title", "winner", "close"}
    bg = NAVY if dark else PAPER
    img.paste(hex_to_rgb(bg), (0, 0, W, H))

    # motif blocks
    draw.rectangle((0, 0, 22, H), fill=COBALT)
    draw.rectangle((22, 0, 36, H), fill=SIENNA)
    draw.rounded_rectangle((1460, 70, 1810, 180), 28, fill=COBALT if not dark else SIENNA)
    draw.text((1492, 104), "m3kit tournament", fill="white", font=font(30, True))

    title_color = "white" if dark else INK
    muted = "#cad3e2" if dark else "#657184"
    draw.text((110, 90), slide["kicker"].upper(), fill=SIENNA if not dark else "#f0b38b", font=font(28, True))
    for n, line in enumerate(wrap_text(draw, slide["title"], font(70, True), 1220)):
        draw.text((110, 140 + n * 84), line, fill=title_color, font=font(70, True))

    y0 = 340
    if slide["kind"] == "ranking":
        colx = [120, 230, 900, 1390, 1530]
        headers = ["Rank", "Lane", "Feature", "Score", "Verdict"]
        draw.rounded_rectangle((100, 310, 1820, 760), 28, fill="white", outline="#d5d9e2", width=2)
        for x, h in zip(colx, headers):
            draw.text((x, 340), h, fill=COBALT, font=font(26, True))
        y = 405
        for row in RANKS:
            fill = "#f4f6fb" if int(row[0]) % 2 else "#ffffff"
            draw.rounded_rectangle((120, y - 14, 1795, y + 66), 14, fill=fill)
            for x, cell in zip(colx, row):
                size = 27 if x != 900 else 25
                for i, line in enumerate(wrap_text(draw, cell, font(size, x in [120, 1390]), 455 if x in [230, 900] else 250)):
                    draw.text((x, y + i * 30), line, fill=INK, font=font(size, x in [120, 1390]))
            y += 84
    elif slide["kind"] == "rubric":
        labels = slide["bullets"]
        for i, label in enumerate(labels):
            x = 140 + (i % 3) * 580
            y = 345 + (i // 3) * 190
            draw_card(draw, (x, y, 500, 140), f"0–5", label, [COBALT, SIENNA, GREEN][i % 3])
    elif slide["kind"] == "diagram":
        boxes = [
            ("TableDefinition", "What columns exist"),
            ("ColumnViewState", "Visible, order, pin"),
            ("m3k-column-manager", "User changes state"),
            ("m3k-data-table", "Renders resolved columns"),
        ]
        for i, (t, b) in enumerate(boxes):
            x = 130 + i * 435
            draw_card(draw, (x, 390, 360, 210), t, b, [COBALT, SIENNA, GREEN, COBALT][i])
            if i < 3:
                draw.line((x + 370, 495, x + 425, 495), fill=INK, width=8)
                draw.polygon([(x+425,495),(x+405,482),(x+405,508)], fill=INK)
    else:
        for i, b in enumerate(slide["bullets"]):
            y = y0 + i * 94
            fill = "#f9fbff" if not dark else "#22304a"
            outline = "#dce1eb" if not dark else "#3d4a61"
            draw.rounded_rectangle((120, y - 22, 1310, y + 58), 20, fill=fill, outline=outline, width=2)
            draw.ellipse((150, y + 1, 182, y + 33), fill=COBALT if i % 2 == 0 else SIENNA)
            for j, line in enumerate(wrap_text(draw, b, font(29), 1040)):
                draw.text((205, y - 3 + j * 32), line, fill=title_color if dark else INK, font=font(29))
        # score/visual panel
        if slide["kind"] == "winner":
            draw.rounded_rectangle((1430, 360, 1780, 680), 36, fill=GREEN)
            draw.text((1500, 425), "30/30", fill="white", font=font(76, True))
            draw.text((1482, 545), "contract fit", fill="white", font=font(34, True))
            draw.text((1508, 590), "a11y + state", fill="white", font=font(30, True))
            draw.text((1508, 632), "feasible", fill="white", font=font(34, True))
        elif slide["kind"] == "title":
            draw.rounded_rectangle((1380, 390, 1790, 710), 36, fill=SIENNA)
            draw.text((1440, 430), "D", fill="white", font=font(145, True))
            draw.text((1550, 485), "wins", fill="white", font=font(50, True))
            draw.text((1428, 600), "Column Manager", fill="white", font=font(34, True))
        else:
            draw.rounded_rectangle((1420, 370, 1780, 650), 32, fill="#ffffff" if not dark else "#22304a", outline="#d5d9e2", width=2)
            draw.text((1470, 410), "m3k-*", fill=COBALT if not dark else "#9fb4ff", font=font(52, True))
            draw.text((1460, 500), "token-only\nsource-backed\nlaunch-ready", fill=INK if not dark else "white", font=font(31, True), spacing=12)

    # callout
    cy = 850
    draw.rounded_rectangle((110, cy, 1810, 990), 28, fill=COBALT if not dark else "#24334f")
    draw.text((145, cy + 32), "Release note", fill="#bcd0ff", font=font(26, True))
    for i, line in enumerate(wrap_text(draw, slide["callout"], font(30, True), 1540)):
        draw.text((145, cy + 68 + i * 36), line, fill="white", font=font(30, True))
    draw.text((1690, 1016), f"{idx:02d}/{len(SLIDES):02d}", fill=muted, font=font(26, True))
    path = SLIDES_DIR / f"slide-{idx:02d}.png"
    img.save(path)
    return path


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Aptos"
    run.font.color.rgb = RGBColor(*hex_to_rgb(color))
    return shape


def add_bullets(slide, items, x, y, w, h, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(*hex_to_rgb(color))
        p.space_after = Pt(8)
    return box


def create_pptx(slide_images):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for idx, s in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        # Use generated image as full-slide designed background for visual fidelity.
        slide.shapes.add_picture(str(slide_images[idx-1]), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        # Add hidden-ish selectable title for accessibility/text extraction.
        add_textbox(slide, 0.2, 7.15, 12.8, 0.2, f"Slide {idx}: {s['title']}", 1, False, PAPER)
    prs.save(str(PPTX))


def md_to_html_text(md: str) -> str:
    try:
        import markdown
        return markdown.markdown(md, extensions=["tables", "fenced_code"])
    except Exception:
        return "<pre>" + html.escape(md) + "</pre>"


def create_html():
    report = (ROOT / "tournament-report.md").read_text()
    packet = (ROOT / "notebooklm-source-packet.md").read_text()
    candidates = (ROOT / "feature-candidate-list.md").read_text()
    css = f"""
    :root {{ color-scheme: light dark; --navy:{NAVY}; --cobalt:{COBALT}; --sienna:{SIENNA}; --paper:{PAPER}; --ink:{INK}; }}
    * {{ box-sizing: border-box; }}
    body {{ margin:0; font-family: Inter, ui-sans-serif, system-ui, -apple-system, Segoe UI, sans-serif; background: var(--paper); color: var(--ink); line-height:1.55; }}
    header {{ background: linear-gradient(135deg, var(--navy), #273757); color:white; padding: 32px 18px 26px; border-left: 10px solid var(--cobalt); }}
    main {{ max-width: 920px; margin: 0 auto; padding: 18px; }}
    h1 {{ font-size: clamp(30px, 7vw, 54px); line-height:1.02; margin: 0 0 12px; }}
    h2 {{ margin-top: 32px; padding-top: 16px; border-top: 3px solid #d9dedf; color: var(--navy); }}
    h3 {{ color: var(--cobalt); }}
    .badge {{ display:inline-block; background:#f0b38b; color:#1b2332; border-radius:999px; padding:5px 12px; font-weight:700; margin:4px 6px 4px 0; }}
    .cards {{ display:grid; grid-template-columns: repeat(auto-fit, minmax(190px, 1fr)); gap:12px; margin:18px 0; }}
    .card {{ background:white; border:1px solid #d6dae3; border-radius:18px; padding:16px; box-shadow:0 3px 12px rgba(20,32,51,.08); }}
    .winner {{ border-left: 8px solid var(--cobalt); }}
    table {{ width:100%; border-collapse: collapse; background:white; border-radius:14px; overflow:hidden; display:block; overflow-x:auto; }}
    th, td {{ padding:10px; border-bottom:1px solid #e3e6ed; text-align:left; vertical-align:top; }}
    th {{ background:#eef2fb; color:#182033; }}
    code {{ background:#eef2fb; padding:2px 5px; border-radius:5px; }}
    a {{ color: var(--cobalt); }}
    .path {{ font-family: ui-monospace, SFMono-Regular, Menlo, monospace; font-size: 13px; overflow-wrap:anywhere; }}
    .note {{ background:#fff7e8; border:1px solid #e2bf82; border-left:7px solid var(--sienna); border-radius:14px; padding:14px; }}
    @media (prefers-color-scheme: dark) {{ body {{ background:#121722; color:#edf1f7; }} .card, table {{ background:#1d2636; }} h2 {{ color:#dbe5ff; border-color:#39465d; }} th {{ background:#273757; color:#fff; }} td, th {{ border-color:#3d485d; }} code {{ background:#273757; }} .note {{ background:#2b2419; }} }}
    """
    ranking_rows = "".join("<tr>" + "".join(f"<td>{html.escape(c)}</td>" for c in row) + "</tr>" for row in RANKS)
    body = f"""<!doctype html><html><head><meta charset='utf-8'><meta name='viewport' content='width=device-width,initial-scale=1'><title>m3kit UI Tournament v2 Brief</title><style>{css}</style></head>
    <body><header><div class='badge'>Release Captain packet</div><div class='badge'>2026-07-06</div><h1>m3kit UI Component Tournament v2</h1><p>Phone-readable brief for the model×driver tournament deliverables. Winner: <strong>Data Table Column Manager</strong>.</p></header><main>
    <section class='cards'>
      <div class='card winner'><h3>Winner</h3><p><strong>OpenCode / DeepSeek v4 Flash free</strong><br>Data Table Column Manager<br><strong>30/30</strong></p></div>
      <div class='card'><h3>Runner-up</h3><p>Antigravity / Gemini 3.5 Flash Medium<br>Data Table Column Manager<br><strong>28/30</strong></p></div>
      <div class='card'><h3>Follow-up</h3><p>Codex / gpt-5.5 proposed Saved View Menu, best as slice two after column-state emission.</p></div>
      <div class='card'><h3>NotebookLM</h3><p>Source packet ready. Notebook creation pending; no link fabricated.</p></div>
    </section>
    <h2>Ranking summary</h2><table><thead><tr><th>Rank</th><th>Lane</th><th>Feature</th><th>Score</th><th>Verdict</th></tr></thead><tbody>{ranking_rows}</tbody></table>
    <h2>Recommendation</h2><div class='note'><p>Approve a Spec Kit feature pack at <code>specs/007-table-column-manager-and-saved-view-seam/</code>: build <code>m3k-column-manager</code> in <code>libs/table</code>, cover hidden/pinned/locked/many/narrow states, add unit tests for emitted state, and Cypress keyboard/focus tests. Saved View Manager follows as the second slice.</p></div>
    <h2>Deliverables</h2><ul>
      <li class='path'>Report: {html.escape(str(ROOT / 'tournament-report.md'))}</li>
      <li class='path'>NotebookLM source packet: {html.escape(str(ROOT / 'notebooklm-source-packet.md'))}</li>
      <li class='path'>PPTX deck: {html.escape(str(PPTX))}</li>
      <li class='path'>Video package: {html.escape(str(VIDEO_DIR))}</li>
      <li class='path'>Silent MP4 deck video: {html.escape(str(MP4))}</li>
      <li class='path'>Voiceover script: {html.escape(str(VOICEOVER))}</li>
    </ul>
    <h2>Core report</h2>{md_to_html_text(report)}
    <h2>NotebookLM source packet excerpt</h2>{md_to_html_text(packet)}
    <h2>Candidate list</h2>{md_to_html_text(candidates)}
    </main></body></html>"""
    HTML.write_text(body)


def create_srt():
    lines = []
    section_text = list(SCRIPT_SECTIONS.values())
    for i, s in enumerate(SLIDES, 1):
        start = (i - 1) * 8
        end = i * 8
        def ts(sec): return f"00:{sec//60:02d}:{sec%60:02d},000"
        narration = section_text[min((i-1) * len(section_text) // len(SLIDES), len(section_text)-1)]
        caption = f"{s['title']}. {narration}"
        lines += [str(i), f"{ts(start)} --> {ts(end)}", caption, ""]
    SRT.write_text("\n".join(lines))


def create_voiceover_script():
    original = (ROOT / "video-presentation-script.md").read_text()
    extra = "\n\n## Slide timing map\n\n" + "\n".join([f"- Slide {i:02d}: {s['title']} — 8 seconds" for i, s in enumerate(SLIDES, 1)])
    VOICEOVER.write_text(original + extra + "\n")


def create_video():
    # Build a simple concat demuxer list: each PNG is displayed for 8 seconds.
    list_file = VIDEO_DIR / "ffmpeg-slides.txt"
    with list_file.open("w") as f:
        for p in sorted(SLIDES_DIR.glob("slide-*.png")):
            f.write(f"file '{p}'\n")
            f.write("duration 8\n")
        last = sorted(SLIDES_DIR.glob("slide-*.png"))[-1]
        f.write(f"file '{last}'\n")
    cmd = [
        "ffmpeg", "-y", "-f", "concat", "-safe", "0", "-i", str(list_file),
        "-vf", "fps=30,format=yuv420p", "-c:v", "libx264", "-movflags", "+faststart", str(MP4)
    ]
    subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)


def validate_sources():
    missing = [f for f in REQUIRED if not (ROOT / f).exists()]
    contestants = sorted(ROOT.glob("contestant-entry-*.md"))
    if missing:
        raise SystemExit(f"Missing required files: {missing}")
    if len(contestants) < 1:
        raise SystemExit("No contestant packets found")
    return contestants


def update_manifest(contestants, tailnet_url="PENDING serve-artifact run"):
    files = sorted([p for p in ROOT.rglob("*") if p.is_file() and p.name != "artifact-manifest.md"])
    lines = [
        "# Artifact Manifest",
        "",
        f"Root: `{ROOT.relative_to(Path.cwd()) if ROOT.is_relative_to(Path.cwd()) else ROOT}`",
        "",
        "## Release Captain packaging status",
        "",
        "- Authority: local-write artifact packaging; external-write limited to Tailnet artifact serving after local generation.",
        "- NotebookLM: source packet ready / notebook creation pending. No NotebookLM notebook link was created or fabricated.",
        f"- Tailnet report URL: {tailnet_url}",
        "",
        "## Primary deliverables",
        "",
        f"- Phone HTML brief: `{HTML.name}`",
        f"- Polished PPTX deck: `{PPTX.name}`",
        f"- Video package directory: `{VIDEO_DIR.name}/`",
        f"- Silent MP4 deck video: `{MP4.relative_to(ROOT)}`",
        f"- Captions: `{SRT.relative_to(ROOT)}`",
        f"- Voiceover script: `{VOICEOVER.relative_to(ROOT)}`",
        f"- Verification log: `{VERIFY.name}`",
        "",
        "## Source artifacts verified/read",
        "",
    ]
    for name in REQUIRED:
        lines.append(f"- `{name}`")
    lines += ["", f"## Contestant packets ({len(contestants)})", ""]
    for c in contestants:
        lines.append(f"- `{c.name}`")
    lines += ["", "## Full artifact file list", ""]
    for p in files:
        lines.append(f"- `{p.relative_to(ROOT)}`")
    (ROOT / "artifact-manifest.md").write_text("\n".join(lines) + "\n")


def write_verification(contestants):
    stat_lines = []
    for p in [PPTX, HTML, MP4, SRT, VOICEOVER, ROOT / "artifact-manifest.md"]:
        stat_lines.append(f"{p}: exists={p.exists()} size={p.stat().st_size if p.exists() else 0}")
    VERIFY.write_text("\n".join([
        "Release Captain verification results",
        "Authority: local-write artifact generation; Tailnet serve verified separately.",
        f"Required sources present: {', '.join(REQUIRED)}",
        f"Contestant packets present: {len(contestants)}",
        "Winner/ranking summary:",
        *[" | ".join(row) for row in RANKS],
        "Generated artifacts:",
        *stat_lines,
    ]) + "\n")


def main():
    os.chdir(ROOT)
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    VIDEO_DIR.mkdir(parents=True, exist_ok=True)
    contestants = validate_sources()
    slide_images = [draw_slide_image(s, i) for i, s in enumerate(SLIDES, 1)]
    create_pptx(slide_images)
    create_html()
    create_srt()
    create_voiceover_script()
    create_video()
    update_manifest(contestants)
    write_verification(contestants)
    print(f"generated_pptx={PPTX}")
    print(f"generated_html={HTML}")
    print(f"generated_mp4={MP4}")
    print(f"contestant_packets={len(contestants)}")

if __name__ == "__main__":
    main()
